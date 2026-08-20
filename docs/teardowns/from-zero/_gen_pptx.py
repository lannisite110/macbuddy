#!/usr/bin/env python3
"""Generate investigation-report PPT from teardown findings."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Visual direction: slate + teal (avoid purple / cream / newspaper defaults)
BG = RGBColor(0x0F, 0x1C, 0x24)
CARD = RGBColor(0x1A, 0x2C, 0x36)
ACCENT = RGBColor(0x2A, 0x9D, 0x8F)
TEXT = RGBColor(0xF0, 0xF4, 0xF5)
MUTED = RGBColor(0xA8, 0xB5, 0xBC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    p.alignment = align
    return tf


def add_bullets(slide, left, top, width, height, lines, size=16, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
    return tf


def title_bar(slide, title, subtitle=None):
    add_textbox(slide, 0.5, 0.25, 9, 0.5, title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.5, 0.75, 9, 0.4, subtitle, size=14, color=ACCENT)
    # accent line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(2.2), Inches(0.04))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    s = new_slide(prs)
    add_textbox(s, 0.8, 2.0, 11.5, 1.0, "AI 产品拆解调查报告", size=40, bold=True, color=WHITE)
    add_textbox(
        s,
        0.8,
        3.1,
        11.5,
        1.2,
        "Cursor · Codex · WorkBuddy · Qoder · DeepSeek Harness · Unsloth\n基于 docs/teardowns/from-zero 七份材料汇总",
        size=18,
        color=MUTED,
    )
    add_textbox(s, 0.8, 5.8, 11.5, 0.4, "2026-08-20  ·  教学调研用途，非各公司内部文档", size=14, color=ACCENT)

    # 2 Agenda
    s = new_slide(prs)
    title_bar(s, "目录", "本报告与演示结构")
    add_bullets(
        s,
        0.8,
        1.6,
        11,
        5,
        [
            "1. 调查目的与材料范围",
            "2. 核心发现：同一口锅 vs 两类问题",
            "3. 六产品分册要点（含 Unsloth）",
            "4. 横向对比表",
            "5. 对 MacBuddy 的行动建议",
            "6. 结论与阅读顺序",
        ],
        size=20,
    )

    # 3 Purpose
    s = new_slide(prs)
    title_bar(s, "调查目的", "要回答的四个问题")
    add_bullets(
        s,
        0.8,
        1.6,
        11,
        5,
        [
            "各自解决什么痛点？",
            "设计概念与系统积木如何拼起来？",
            "从零开发的最小可行路径是什么？",
            "MacBuddy 该学什么、不该抄什么？",
        ],
        size=22,
    )

    # 4 Materials
    s = new_slide(prs)
    title_bar(s, "材料范围", "七份 from-zero 拆解")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "00 总览 — 共用 Agent 循环；Unsloth 例外",
            "01 Cursor — 把 AI 焊进代码编辑器",
            "02 Codex — Agent 发动机 + 多种外壳",
            "03 WorkBuddy — 办公 Agent，交文档/表格",
            "04 Qoder — 定目标、审结果的编程平台",
            "05 DeepSeek Harness — 一切皆插件的开源底盘",
            "06 Unsloth — 本地训练/跑模型（不是 Agent）",
        ],
        size=18,
    )

    # 5 Same pot
    s = new_slide(prs)
    title_bar(s, "核心发现①：同一口锅", "除 Unsloth 外，内核都是 Agent 循环")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "人说话 → 组装上下文 → 模型决策 → 调用工具 → 观察 → 再决策",
            "差别只在：外壳、工具危险度、上下文从哪来",
            "Cursor：索引 + 编辑器；Codex：App Server 事件",
            "WorkBuddy：资料库 + 产物；Qoder：Quest + 权限",
            "dsh：插件树 + 会话日志即真相",
        ],
        size=18,
    )

    # 6 Two problems
    s = new_slide(prs)
    title_bar(s, "核心发现②：两类问题", "不要把 Unsloth 当成又一个 Agent")
    add_bullets(
        s,
        0.8,
        1.6,
        11.5,
        5,
        [
            "Harness / 产品壳（Cursor / Codex / WorkBuddy / Qoder / dsh）",
            "    → 怎么让模型安全地干活并交付",
            "模型运行时（Unsloth）",
            "    → 怎么让权重在你机器上训得动、跑得动",
            "MacBuddy：薄壳 + sidecar；消费 OpenAI 兼容 /v1 即可",
        ],
        size=18,
    )

    # 7 Cursor
    s = new_slide(prs)
    title_bar(s, "Cursor", "把 AI 焊进编辑器（非插件）")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：毫秒补全、多文件改、仓库检索、可审 diff",
            "概念：Merkle 索引、向量检索、ReAct；Tab/行内/Agent 分快慢",
            "从零：聊天 → Read → ApplyPatch（必预览）→ Grep → 白名单终端",
            "MacBuddy：学补丁预览；不造 Tab 与整套 IDE",
        ],
        size=18,
    )

    # 8 Codex
    s = new_slide(prs)
    title_bar(s, "OpenAI Codex", "发动机 + 标准插座")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：多 UI 不复制循环；批准、取消、沙箱、CI",
            "概念：Thread / Turn / Item；JSON-RPC 可反向问批准",
            "从零：while 循环 → 子进程 JSON 事件 → 批准/取消 → jsonl",
            "MacBuddy：协议写清；SwiftUI 只投影，不写 Agent 状态机",
        ],
        size=18,
    )

    # 9 WorkBuddy
    s = new_slide(prs)
    title_bar(s, "WorkBuddy（腾讯）", "办公工作台：交作业，不只给建议")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：跨应用跳；聊天不交文件；外出遥控电脑",
            "概念：资料库、Skill、连接器、产物篮、云沙箱 vs 本机",
            "从零：授权目录 → 产物盒 → 表格 → 一种办公格式 → 任务状态",
            "MacBuddy：任务有状态、产物一等公民；不抄全家桶",
        ],
        size=18,
    )

    # 10 Qoder
    s = new_slide(prs)
    title_bar(s, "Qoder（阿里）", "Define the goal. Review the result.")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：多文件任务、失忆、想委派整活",
            "概念：四齿轮（Agent/工具/权限/上下文）；Quest；Wiki/Memory",
            "从零：绑仓库 → 只读工具 → diff → 权限三态 → AGENTS.md",
            "MacBuddy：Code 停在补丁+预览；不做完整 IDE / 专家团",
        ],
        size=18,
    )

    # 11 dsh
    s = new_slide(prs)
    title_bar(s, "DeepSeek Harness", "Everything is a plugin · 不是模型")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：核心缠死；多外壳复制；无法回放模型所见",
            "概念：Cordis 插拔；profile/bundle；Turn/Step；日志=真相",
            "铁律：模型可见 ⇒ 必须能从会话日志重建",
            "MacBuddy：只追加事件、sidecar 可替换；不移植 Cordis",
        ],
        size=18,
    )

    # 12 Unsloth
    s = new_slide(prs)
    title_bar(s, "Unsloth", "训练/推理运行时，不是 Agent")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "痛点：太大、全量微调太贵、格式分裂（GGUF/MLX）",
            "NVIDIA：LoRA/QLoRA + Triton 融合内核",
            "Mac：llama.cpp Metal / MLX；与 Ollama 同速度档",
            "MacBuddy：接 /v1；16GB 优先 ≤4B Q4；不自研训练",
        ],
        size=18,
    )

    # 13 Comparison
    s = new_slide(prs)
    title_bar(s, "横向对比（精简）", "主战场与交付物")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "Cursor — IDE 编码 — 代码 diff",
            "Codex — Agent 核心 — 事件流 / 补丁",
            "WorkBuddy — 办公交付 — 文档表格",
            "Qoder — 编程委派 — Quest + diff",
            "dsh — 可插拔底盘 — 会话日志",
            "Unsloth — 训练/推理 — 权重 / GGUF",
        ],
        size=18,
    )

    # 14 Hard decisions
    s = new_slide(prs)
    title_bar(s, "Agent 类共性硬决定", "产品宪法级选择")
    add_bullets(
        s,
        0.8,
        1.6,
        11.5,
        5,
        [
            "改盘前要不要人看 diff / 批准？",
            "工具失败是否把报错原文回灌模型？",
            "界面主线程绝不阻塞等待模型",
            "工作区边界与命令白名单写死",
            "无界面 / CI 模式下 ask 如何降级？",
        ],
        size=20,
    )

    # 15 MacBuddy
    s = new_slide(prs)
    title_bar(s, "对 MacBuddy 的行动建议", "学脊梁，不抄航母")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "保持：薄壳 + sidecar；事件 token/tool/approval/done/error",
            "Work：任务状态 + 产物盒；不做网盘/技能市场",
            "Code：补丁预览 + 权限三态；不做 Quest IDE / Wiki 工厂",
            "模型：本地 Ollama / Unsloth Desktop /v1；小模型保体感",
            "预期：「从零」= 同类最小系统，不是完整商业 IDE",
        ],
        size=17,
    )

    # 16 Conclusion
    s = new_slide(prs)
    title_bar(s, "结论", "值钱的是脊梁，不是皮肤")
    add_bullets(
        s,
        0.8,
        1.6,
        11.5,
        5,
        [
            "值钱：循环 + 工具守卫 + 可审交付 + 可回放上下文",
            "模型怎么跑 与 Agent 怎么干 必须解耦",
            "分册原文：docs/teardowns/from-zero/00–06",
            "完整报告：docs/teardowns/from-zero/调查报告.md",
        ],
        size=20,
    )

    # 17 Reading order
    s = new_slide(prs)
    title_bar(s, "附录：建议阅读顺序", "仍回到七份 md 深读")
    add_bullets(
        s,
        0.8,
        1.5,
        11.5,
        5.5,
        [
            "① 00-README — 同一口锅",
            "② 01-cursor + 02-codex — 循环与外壳",
            "③ 03-workbuddy + 04-qoder — 办公 vs 编程",
            "④ 05-deepseek-harness — 插件与日志铁律",
            "⑤ 06-unsloth — 运行时，勿与 Agent 混淆",
        ],
        size=20,
    )

    # 18 End
    s = new_slide(prs)
    add_textbox(s, 0.8, 2.8, 11.5, 0.8, "谢谢", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(
        s,
        0.8,
        3.8,
        11.5,
        0.8,
        "材料目录：docs/teardowns/from-zero/\n报告：调查报告.md  ·  本演示：AI产品拆解调查报告.pptx",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    out = "/workspace/docs/teardowns/from-zero/AI产品拆解调查报告.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
